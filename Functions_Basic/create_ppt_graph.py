from pptx import Presentation
from pptx.util import Inches
import tkinter as tk
import pathlib
import tkinter.messagebox
from datetime import datetime
import Functions_GUI.buttons as bt


def create_ppt_graph(path: pathlib.WindowsPath = None, function_index: int = None,
                     x: float = 0.15, y: float = 2, width: float =12.96, height: float = 4):
    """
    This function creates a pptx file with given .png files in the given path.
    
    :param path: Path of the directory which contains the .png files.
    :param function_index: Index of the function which is used to create the ppt.
    :param x: x coordinate of the image in the ppt.
    :param y: y coordinate of the image in the ppt.
    :param width: width of the image in the ppt.
    :param height: height of the image in the ppt.
    :return: None
    """
    ##### PPT PATH GENERATION ####
    status_ppt = 0  # status variable holds the information for if the ppt path needs to be changed.
    if not bt.ppt_path:
        # If we don't have any bt.ppt_path that means we will be using the function first time, thus append the default ppt.
        status_ppt = 1  # active status means that we need to save and use another ppt path.
        if bt.basic_settings[1] == 1:  # If DataType-1 selected
            bt.ppt_path = bt.default_path.joinpath('data/Default_Presentation.pptx')
        elif bt.basic_settings[1] == 2:  # If DataType-2 selected
            bt.ppt_path = bt.default_path.joinpath('data/Default_Presentation.pptx')
        else:
            tk.messagebox.showerror('Error on ppt directory',
                                    "Default PPT Scheme could not be found. Please check it before run.")
            return None

    def move_slide(old_index, new_index):
        xml_slides = prs.slides._sldIdLst
        slides = list(xml_slides)
        xml_slides.remove(slides[old_index])
        xml_slides.insert(new_index, slides[old_index])

    prs = Presentation(bt.ppt_path)  # Create the presentation object with given path
    layout = prs.slide_layouts[8]  # 8 th layout slide in the default slides. It represents 'Empty Slide'

    # Iterate over the .png graphics in the directory
    for file in path.iterdir():
        if file.suffix == '.png':
            slide = prs.slides.add_slide(layout)  # Adding empty slide to presentation.
            slide.shapes.add_picture(str(file), Inches(x), Inches(y), width=Inches(width), height=Inches(height))  # Load the image and add it to the slide

            # Move slide to the corresponding spot
            last_slide = len(prs.slides)
            move_slide(last_slide - 1, bt.slide_positions[function_index])  # Moving slide into specific location.

            # Updating the slide positions for further appending
            for i in range(function_index, len(bt.slide_positions)):
                bt.slide_positions[i] += 1

    if status_ppt:
        # Creates new pptx name and file if create_ppt function is used first time.
        now = str(datetime.now())  # Taking current time.
        now = now[:-7]
        now = now.replace(':', '.')
        bt.ppt_path = bt.save_path.joinpath('Auswertung_' + now + '.pptx')  # New save path for the appended default slide.
        prs.save(bt.ppt_path.as_posix())  # Changes path to string without backslashes.
    else:
        prs.save(bt.ppt_path.as_posix())  # If second use, we append the used slide instead of default slide
