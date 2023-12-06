from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import pandas as pd
import tkinter as tk
import tkinter.messagebox
from datetime import datetime
import Functions_GUI.buttons as bt


def create_ppt_table(function_index: int, table_df: pd.DataFrame = None, background_color: bool = False,
                     labels: list = None, color_indexes: list = None, starting_color_index: int = 0):
    """
    Takes the function type as input. Creates table and appends Default ppt for Bosch or Conti.

    :param table_df: Result dataframe which includes printed table data
    :param function_index: Defines which function is the input for create_ppt_table function.
    Function changes behaviour due to creation of table according to selected function.
    :param background_color: Decides that table will conditionally change background color of cells.
    :param labels: Used for special process for specific columns or labels (ex: coloring)
    :param color_indexes: Used to hold colored page number array in TOC table
    :param starting_color_index: Used to detect correct cell for background color of TOC Table of OBD-Radar
    :return: Updates report ppt itself and bt.ppt_path for further ppt path usage.
    """
    status_ppt = 0  # status variable holds the information for if the ppt path needs to be changed.
    if not bt.ppt_path:
        # If we don't have any bt.ppt_path that means we will be using the function first time, thus append the default ppt.
        status_ppt = 1  # active status means that we need to save and use another ppt path.
        if bt.basic_settings[1] == 1:  # If Bosch selected
            bt.ppt_path = bt.default_path.joinpath('data/Default_Presentation.pptx')
        elif bt.basic_settings[1] == 2:  # If Conti selected
            bt.ppt_path = bt.default_path.joinpath('data/Default_Presentation.pptx')
        else:
            tk.messagebox.showerror('Error on ppt directory', "Default PPT Scheme could not be found. Please check it before run.")
            return None
    row_counter = 0  # External counter to append table with corresponding dataframe cell
    background_row_counter = starting_color_index * 20  # External counter to select correct cell for OBD-Radar TOC Table

    def move_slide(old_index, new_index):
        xml_slides = prs.slides._sldIdLst
        slides = list(xml_slides)
        xml_slides.remove(slides[old_index])
        xml_slides.insert(new_index, slides[old_index])

    prs = Presentation(bt.ppt_path)  # Create the presentation object with given path
    layout = prs.slide_layouts[8]  # 8 th layout slide in the default slides. It represents 'Empty Slide'

    # Deciding how many pages needed
    rows = len(table_df.index)
    cols = len(table_df.columns)
    table_df.fillna('N/A', inplace=True)

    # Deciding font size, column widths, table starting positions, table row amount for different types of outputs
    if function_index == 8:  # Function index for Verwendete_DiagRA output
        column_widths = [Inches(2.75), Inches(1), Inches(1.5), Inches(0.75),
                         Inches(0.9), Inches(0.75), Inches(1.6)]
        header_font_size = 11
        row_font_size = 8
        x, y, cx, cy = Inches(0.3), Inches(0.3), Inches(6), Inches(0.3)  # left, up, width of cell, height of cell
        row_amount = 28  #  Amount of row for function types
        page_amount = rows / row_amount  # Whole page
        leftover = rows % row_amount  # Leftover rows
    elif function_index == 9:  # Function index for P_Code output
        column_widths = [Inches(1.85), Inches(0.6), Inches(0.5), Inches(1.55),
                         Inches(4.25), Inches(0.9), Inches(0.9), Inches(0.8),
                         Inches(0.7), Inches(0.6), Inches(0.4)]
        header_font_size = 8
        row_font_size = 6
        x, y, cx, cy = Inches(0.15), Inches(1.6), Inches(6), Inches(0.35)  # left, up, width of cell, height of cell
        row_amount = 26  # Amount of row for function types
        page_amount = rows / row_amount  # Whole page
        leftover = rows % row_amount  # Leftover rows
    elif function_index == 10:  # Function index for Mode9 output
        column_widths = [Inches(0.9), Inches(1.3), Inches(1.7), Inches(1.3),
                         Inches(1.2), Inches(0.9), Inches(0.9), Inches(0.9),
                         Inches(0.9), Inches(0.9), Inches(0.9), Inches(0.9)]
        header_font_size = 8
        row_font_size = 7
        x, y, cx, cy = Inches(0.15), Inches(1.6), Inches(6), Inches(0.35)  # left, up, width of cell, height of cell
        row_amount = 25  # Amount of row for function types
        page_amount = rows / row_amount  # Whole page
        leftover = rows % row_amount  # Leftover rows
    elif function_index == 1:  # Function index for Mode9_Aktuellste output
        column_widths = [Inches(1.7)] + [Inches(1.5)] * (cols - 1)
        header_font_size = 8
        row_font_size = 9
        x, y, cx, cy = Inches(0.15), Inches(1.85), Inches(6), Inches(0.35)  # left, up, width of cell, height of cell
        row_amount = rows  # Amount of row for function types
        page_amount = rows / row_amount  # Whole page
        leftover = 0  # Leftover rows
    elif function_index == 11:  # Function index for IUMPR output
        column_widths = [Inches(3.25)] + [Inches(1.5)] * 6
        header_font_size = 8
        row_font_size = 7
        x, y, cx, cy = Inches(0.15), Inches(1.6), Inches(6), Inches(0.35)  # left, up, width of cell, height of cell
        row_amount = 25  # Amount of row for function types
        page_amount = rows / row_amount  # Whole page
        leftover = rows % row_amount  # Leftover rows
    elif function_index == 3:  # Function index for IUMPR_Aktuellste output
        column_widths = [Inches(4)] + [Inches(1.5)] * (cols - 1)
        header_font_size = 8
        row_font_size = 7
        x, y, cx, cy = Inches(0.15), Inches(1.6), Inches(6), Inches(0.35)  # left, up, width of cell, height of cell
        row_amount = 15  # Amount of row for function types
        page_amount = rows / row_amount  # Whole page
        leftover = rows % row_amount  # Leftover rows
    elif function_index == 2:  # Function index for Mode9 TOC output
        column_widths = [Inches(4.25), Inches(0.7), Inches(4.25), Inches(0.7)]
        header_font_size = 10
        row_font_size = 9
        x, y, cx, cy = Inches(0.15), Inches(1.6), Inches(6), Inches(0.35)  # left, up, width of cell, height of cell
        row_amount = 10  # Amount of row for function types
        page_amount = rows / row_amount  # Whole page
        leftover = rows % row_amount  # Leftover rows
    elif function_index == 4:  # Function index for IUMPR TOC output
        column_widths = [Inches(4.25), Inches(0.7), Inches(4.25), Inches(0.7)]
        header_font_size = 10
        row_font_size = 9
        x, y, cx, cy = Inches(0.15), Inches(1.6), Inches(6), Inches(0.35)  # left, up, width of cell, height of cell
        row_amount = 10  # Amount of row for function types
        page_amount = rows / row_amount  # Whole page
        leftover = rows % row_amount  # Leftover rows
    elif function_index == 5:  # Function index for Ramzellen TOC output
        column_widths = [Inches(4.25), Inches(0.7), Inches(4.25), Inches(0.7)]
        header_font_size = 10
        row_font_size = 9
        x, y, cx, cy = Inches(0.15), Inches(1.6), Inches(6), Inches(0.35)  # left, up, width of cell, height of cell
        row_amount = 10  # Amount of row for function types
        page_amount = rows / row_amount  # Whole page
        leftover = rows % row_amount  # Leftover rows
    elif function_index == 6:  # Function index for OBD-Radar TOC output
        column_widths = [Inches(5.75), Inches(0.7), Inches(5.75), Inches(0.7)]
        header_font_size = 10
        row_font_size = 8
        x, y, cx, cy = Inches(0.15), Inches(1.6), Inches(6), Inches(0.35)  # left, up, width of cell, height of cell
        row_amount = 10  # Amount of row for function types
        page_amount = rows / row_amount  # Whole page
        leftover = rows % row_amount  # Leftover rows
    else:
        column_widths = []
        header_font_size = 1
        row_font_size = 1
        x, y, cx, cy = Inches(0.3), Inches(0.3), Inches(6), Inches(0.3)  # left, up, width of cell, height of cell
        row_amount = 0  # Amount of row for function types
        page_amount = 0  # Whole page
        leftover = 0  # Leftover rows

    # Appending ppt
    while page_amount >= 0:
        slide = prs.slides.add_slide(layout)  # Adding empty slide to presentation.

        # Add table to the slide
        if page_amount < 1 and leftover == 0:  # Means that there is no leftover row left. No page will be created
            break
        elif page_amount < 1 and leftover > 0:  # Means added page is leftover page which has leftover + 1 (header) rows
            table = slide.shapes.add_table(leftover + 1, cols, x, y, cx, cy).table
        else:  # Means added page is whole page, +1 for header
            table = slide.shapes.add_table(row_amount + 1, cols, x, y, cx, cy).table

        # Table style
        table.style = "Table Grid"

        # Assign column widths
        if column_widths:
            for i in range(cols):
                table.columns[i].width = column_widths[i]
        else:
            print('Error on function index, table columns could not be assigned!')
            return 0

        # Assigning header names
        if function_index == 1 or function_index == 3 or function_index == 11:
            table_df.columns = table_df.columns.astype(str)

        for i in range(cols):
            table.cell(0, i).text = table_df.columns[i]
            table.cell(0, i).text_frame.paragraphs[0].font.size = Pt(header_font_size)
            table.cell(0, i).text_frame.paragraphs[0].font.name = 'Calibri'
            table.cell(0, i).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Filling the table
        if page_amount < 1:  # Leftover page
            for i in range(leftover):
                for j in range(cols):
                    table.cell(i + 1, j).text = table_df.iloc[row_counter + i, j]
                    table.cell(i + 1, j).text_frame.paragraphs[0].font.size = Pt(row_font_size)
                    table.cell(i + 1, j).text_frame.paragraphs[0].font.name = 'Calibri'
        else:
            for i in range(row_amount):  # Full page
                for j in range(cols):
                    table.cell(i + 1, j).text = table_df.iloc[row_counter + i, j]
                    table.cell(i + 1, j).text_frame.paragraphs[0].font.size = Pt(row_font_size)
                    table.cell(i + 1, j).text_frame.paragraphs[0].font.name = 'Calibri'

        # Conditional background coloring
        if background_color:
            if function_index == 10:  # Mode9 Function
                if page_amount < 1:  # Leftover page
                    for i in range(leftover):
                        for j in range(cols):
                            if j >= 6:  # Colored columns
                                try:
                                    ratio = float(table_df.iloc[row_counter + i, j].split()[0])
                                    nominator_denominator = table_df.iloc[row_counter + i, j].split()[1]
                                except ValueError:
                                    continue
                                if '0/0' in nominator_denominator:
                                    continue
                                elif ratio >= 0.5:
                                    table.cell(i + 1, j).fill.solid()
                                    table.cell(i + 1, j).fill.fore_color.rgb = RGBColor(0, 255, 0)  # Set background color to green
                                elif ratio >= 0.334:
                                    table.cell(i + 1, j).fill.solid()
                                    table.cell(i + 1, j).fill.fore_color.rgb = RGBColor(255, 255, 0)  # Set background color to orange
                                else:
                                    table.cell(i + 1, j).fill.solid()
                                    table.cell(i + 1, j).fill.fore_color.rgb = RGBColor(255, 0, 0)  # Set background color to red
                else:
                    for i in range(row_amount):
                        for j in range(cols):
                            if j >= 6:  # Colored columns
                                try:
                                    ratio = float(table_df.iloc[row_counter + i, j].split()[0])
                                    nominator_denominator = table_df.iloc[row_counter + i, j].split()[1]
                                except ValueError:
                                    continue
                                if '0/0' in nominator_denominator:
                                    continue
                                elif ratio >= 0.5:
                                    table.cell(i + 1, j).fill.solid()
                                    table.cell(i + 1, j).fill.fore_color.rgb = RGBColor(0, 255,0)  # Set background color to green
                                elif ratio >= 0.334:
                                    table.cell(i + 1, j).fill.solid()
                                    table.cell(i + 1, j).fill.fore_color.rgb = RGBColor(255, 255, 0)  # Set background color to orange
                                else:
                                    table.cell(i + 1, j).fill.solid()
                                    table.cell(i + 1, j).fill.fore_color.rgb = RGBColor(255, 0, 0)  # Set background color to red
            elif function_index == 1:  # Mode9 Aktuellste Function
                last_colored_column_index = cols - 1
                check_flotte = False
                if page_amount < 1:  # Leftover page
                    for i in range(leftover):
                        if i >= 5:
                            for j in range(cols):
                                if 1 <= j <= last_colored_column_index:
                                    nominator_denominator = '1/1'  # Default nominator/denominator value
                                    check_flotte = False
                                    if table_df.columns.tolist()[j] == ' ':  # Flotte column detection
                                        try:
                                            cell_value = table_df.iloc[row_counter + i, j].split()[0]
                                            if cell_value == '100':
                                                # Detects if flotte calculated or given directly. 100 -> given 100.0 calculated
                                                continue
                                            ratio = float(cell_value)
                                            check_flotte = True
                                        except ValueError:
                                            continue
                                    elif table_df.columns.tolist()[j] != '':  # Durchschnitt column detection
                                        try:
                                            ratio = float(table_df.iloc[row_counter + i, j].split()[0])
                                            nominator_denominator = table_df.iloc[row_counter + i, j].split()[1]
                                        except ValueError:
                                            continue
                                    else:
                                        try:
                                            ratio = float(table_df.iloc[row_counter + i, j].split()[0])
                                            nominator_denominator = table_df.iloc[row_counter + i, j].split()[1]
                                        except ValueError:
                                            continue

                                    if '0/0' in nominator_denominator:
                                        continue
                                    elif check_flotte:
                                        if ratio >= 50:
                                            table.cell(i + 1, j).fill.solid()
                                            table.cell(i + 1, j).fill.fore_color.rgb = RGBColor(0, 255, 0)  # Set background color to green
                                        elif ratio >= 34:
                                            table.cell(i + 1, j).fill.solid()
                                            table.cell(i + 1, j).fill.fore_color.rgb = RGBColor(255, 255, 0)  # Set background color to orange
                                        else:
                                            table.cell(i + 1, j).fill.solid()
                                            table.cell(i + 1, j).fill.fore_color.rgb = RGBColor(255, 0, 0)  # Set background color to red
                                    elif ratio >= 0.5:
                                        table.cell(i + 1, j).fill.solid()
                                        table.cell(i + 1, j).fill.fore_color.rgb = RGBColor(0, 255, 0)  # Set background color to green
                                    elif ratio >= 0.334:
                                        table.cell(i + 1, j).fill.solid()
                                        table.cell(i + 1, j).fill.fore_color.rgb = RGBColor(255, 255, 0)  # Set background color to orange
                                    else:
                                        table.cell(i + 1, j).fill.solid()
                                        table.cell(i + 1, j).fill.fore_color.rgb = RGBColor(255, 0, 0)  # Set background color to red
                else:
                    for i in range(row_amount):
                        if i >= 5:
                            for j in range(cols):
                                if 1 <= j <= last_colored_column_index:
                                    nominator_denominator = '1/1'  # Default nominator/denominator value
                                    check_flotte = False
                                    if table_df.columns.tolist()[j] == ' ':  # Flotte column detection
                                        try:
                                            cell_value = table_df.iloc[row_counter + i, j].split()[0]
                                            if cell_value == '100':
                                                # Detects if flotte calculated or given directly. 100 -> given 100.0 calculated
                                                continue
                                            ratio = float(cell_value)
                                            check_flotte = True
                                        except ValueError:
                                            continue
                                    elif table_df.columns.tolist()[j] != '':  # Durchschnitt column detection
                                        try:
                                            ratio = float(table_df.iloc[row_counter + i, j].split()[0])
                                            nominator_denominator = table_df.iloc[row_counter + i, j].split()[1]
                                        except ValueError:
                                            continue
                                    else:
                                        try:
                                            ratio = float(table_df.iloc[row_counter + i, j].split()[0])
                                            nominator_denominator = table_df.iloc[row_counter + i, j].split()[1]
                                        except ValueError:
                                            continue
                                    if '0/0' in nominator_denominator:
                                        continue
                                    elif check_flotte:
                                        if ratio >= 50:
                                            table.cell(i + 1, j).fill.solid()
                                            table.cell(i + 1, j).fill.fore_color.rgb = RGBColor(0, 255, 0)  # Set background color to green
                                        elif ratio >= 34:
                                            table.cell(i + 1, j).fill.solid()
                                            table.cell(i + 1, j).fill.fore_color.rgb = RGBColor(255, 255, 0)  # Set background color to orange
                                        else:
                                            table.cell(i + 1, j).fill.solid()
                                            table.cell(i + 1, j).fill.fore_color.rgb = RGBColor(255, 0, 0)  # Set background color to red
                                    elif ratio >= 0.5:
                                        table.cell(i + 1, j).fill.solid()
                                        table.cell(i + 1, j).fill.fore_color.rgb = RGBColor(0, 255, 0)  # Set background color to green
                                    elif ratio >= 0.334:
                                        table.cell(i + 1, j).fill.solid()
                                        table.cell(i + 1, j).fill.fore_color.rgb = RGBColor(255, 255, 0)  # Set background color to orange
                                    else:
                                        table.cell(i + 1, j).fill.solid()
                                        table.cell(i + 1, j).fill.fore_color.rgb = RGBColor(255, 0, 0)  # Set background color to red
            elif function_index == 11:  # IUMPR Function
                if page_amount < 1:  # Leftover page
                    for i in range(leftover):
                        if table_df.iloc[row_counter + i, 0] in labels:
                            for j in range(cols):
                                if j >= 1:  # Colored columns
                                    try:
                                        ratio = float(table_df.iloc[row_counter + i, j].split()[0])
                                        nominator_denominator = table_df.iloc[row_counter + i, j].split()[1]
                                    except ValueError:
                                        continue
                                    if '0/0' in nominator_denominator:
                                        continue
                                    elif ratio >= 0.5:
                                        table.cell(i + 1, j).fill.solid()
                                        table.cell(i + 1, j).fill.fore_color.rgb = RGBColor(0, 255, 0)  # Set background color to green
                                    elif ratio >= 0.334:
                                        table.cell(i + 1, j).fill.solid()
                                        table.cell(i + 1, j).fill.fore_color.rgb = RGBColor(255, 255, 0)  # Set background color to orange
                                    else:
                                        table.cell(i + 1, j).fill.solid()
                                        table.cell(i + 1, j).fill.fore_color.rgb = RGBColor(255, 0, 0)  # Set background color to red
                else:
                    for i in range(row_amount):
                        if table_df.iloc[row_counter + i, 0] in labels:
                            for j in range(cols):
                                if j >= 1:  # Colored columns
                                    try:
                                        ratio = float(table_df.iloc[row_counter + i, j].split()[0])
                                        nominator_denominator = table_df.iloc[row_counter + i, j].split()[1]
                                    except ValueError:
                                        continue
                                    if '0/0' in nominator_denominator:
                                        continue
                                    elif ratio >= 0.5:
                                        table.cell(i + 1, j).fill.solid()
                                        table.cell(i + 1, j).fill.fore_color.rgb = RGBColor(0, 255,0)  # Set background color to green
                                    elif ratio >= 0.334:
                                        table.cell(i + 1, j).fill.solid()
                                        table.cell(i + 1, j).fill.fore_color.rgb = RGBColor(255, 255, 0)  # Set background color to orange
                                    else:
                                        table.cell(i + 1, j).fill.solid()
                                        table.cell(i + 1, j).fill.fore_color.rgb = RGBColor(255, 0, 0)  # Set background color to red
            elif function_index == 3:  # IUMPR Aktuellste Function
                last_colored_column_index = cols - 1
                if page_amount < 1:  # Leftover page
                    for i in range(leftover):
                        if table_df.iloc[row_counter + i, 0] in labels:
                            for j in range(cols):
                                if 1 <= j <= last_colored_column_index:  # Colored columns
                                    nominator_denominator = '1/1'  # Default nominator/denominator value
                                    check_flotte = False
                                    if table_df.columns.tolist()[j] == ' ':  # Flotte column detection
                                        try:
                                            cell_value = table_df.iloc[row_counter + i, j].split()[0]
                                            if cell_value == '100':
                                                # Detects if flotte calculated or given directly. 100 -> given 100.0 calculated
                                                continue
                                            ratio = float(cell_value)
                                            check_flotte = True
                                        except ValueError:
                                            continue
                                    elif table_df.columns.tolist()[j] != '':  # Durchschnitt column detection
                                        try:
                                            ratio = float(table_df.iloc[row_counter + i, j].split()[0])
                                            nominator_denominator = table_df.iloc[row_counter + i, j].split()[1]
                                        except ValueError:
                                            continue
                                    else:  # Ratio values column
                                        try:
                                            ratio = float(table_df.iloc[row_counter + i, j].split()[0])
                                            nominator_denominator = table_df.iloc[row_counter + i, j].split()[1]
                                        except ValueError:
                                            continue

                                    if '0/0' in nominator_denominator:
                                        continue
                                    elif check_flotte:
                                        if ratio >= 50:
                                            table.cell(i + 1, j).fill.solid()
                                            table.cell(i + 1, j).fill.fore_color.rgb = RGBColor(0, 255, 0)  # Set background color to green
                                        elif ratio >= 34:
                                            table.cell(i + 1, j).fill.solid()
                                            table.cell(i + 1, j).fill.fore_color.rgb = RGBColor(255, 255, 0)  # Set background color to orange
                                        else:
                                            table.cell(i + 1, j).fill.solid()
                                            table.cell(i + 1, j).fill.fore_color.rgb = RGBColor(255, 0, 0)  # Set background color to red
                                    elif ratio >= 0.5:
                                        table.cell(i + 1, j).fill.solid()
                                        table.cell(i + 1, j).fill.fore_color.rgb = RGBColor(0, 255, 0)  # Set background color to green
                                    elif ratio >= 0.334:
                                        table.cell(i + 1, j).fill.solid()
                                        table.cell(i + 1, j).fill.fore_color.rgb = RGBColor(255, 255, 0)  # Set background color to orange
                                    else:
                                        table.cell(i + 1, j).fill.solid()
                                        table.cell(i + 1, j).fill.fore_color.rgb = RGBColor(255, 0, 0)  # Set background color to red
                else:
                    for i in range(row_amount):
                        if table_df.iloc[row_counter + i, 0] in labels:
                            for j in range(cols):
                                if 1 <= j <= last_colored_column_index:  # Colored columns
                                    nominator_denominator = '1/1'  # Default nominator/denominator value
                                    check_flotte = False
                                    if table_df.columns.tolist()[j] == ' ':  # Flotte column detection
                                        try:
                                            cell_value = table_df.iloc[row_counter + i, j].split()[0]
                                            if cell_value == '100':
                                                # Detects if flotte calculated or given directly. 100 -> given 100.0 calculated
                                                continue
                                            ratio = float(cell_value)
                                            check_flotte = True
                                        except ValueError:
                                            continue
                                    elif table_df.columns.tolist()[j] != '':  # Durchschnitt column detection
                                        try:
                                            ratio = float(table_df.iloc[row_counter + i, j].split()[0])
                                            nominator_denominator = table_df.iloc[row_counter + i, j].split()[1]
                                        except ValueError:
                                            continue
                                    else:  # Ratio values column
                                        try:
                                            ratio = float(table_df.iloc[row_counter + i, j].split()[0])
                                            nominator_denominator = table_df.iloc[row_counter + i, j].split()[1]
                                        except ValueError:
                                            continue

                                    if '0/0' in nominator_denominator:
                                        continue
                                    elif check_flotte:
                                        if ratio >= 50:
                                            table.cell(i + 1, j).fill.solid()
                                            table.cell(i + 1, j).fill.fore_color.rgb = RGBColor(0, 255, 0)  # Set background color to green
                                        elif ratio >= 34:
                                            table.cell(i + 1, j).fill.solid()
                                            table.cell(i + 1, j).fill.fore_color.rgb = RGBColor(255, 255, 0)  # Set background color to orange
                                        else:
                                            table.cell(i + 1, j).fill.solid()
                                            table.cell(i + 1, j).fill.fore_color.rgb = RGBColor(255, 0, 0)  # Set background color to red
                                    elif ratio >= 0.5:
                                        table.cell(i + 1, j).fill.solid()
                                        table.cell(i + 1, j).fill.fore_color.rgb = RGBColor(0, 255, 0)  # Set background color to green
                                    elif ratio >= 0.334:
                                        table.cell(i + 1, j).fill.solid()
                                        table.cell(i + 1, j).fill.fore_color.rgb = RGBColor(255, 255, 0)  # Set background color to orange
                                    else:
                                        table.cell(i + 1, j).fill.solid()
                                        table.cell(i + 1, j).fill.fore_color.rgb = RGBColor(255, 0, 0)  # Set background color to red
            elif function_index == 6:  # OBD-Radar TOC Table
                if page_amount < 1:  # Leftover page
                    for j in range(cols):
                        if j == 1 or j == 3:  # Colored columns
                            for i in range(leftover):
                                if table_df.iloc[i, j]:
                                    color = color_indexes[background_row_counter]

                                    if color == 'red':
                                        table.cell(i + 1, j).fill.solid()
                                        table.cell(i + 1, j).fill.fore_color.rgb = RGBColor(255, 0, 0)  # Set background color to red
                                    elif color == 'yellow':
                                        table.cell(i + 1, j).fill.solid()
                                        table.cell(i + 1, j).fill.fore_color.rgb = RGBColor(255, 255,0)  # Set background color to yellow
                                    elif color == 'green':
                                        table.cell(i + 1, j).fill.solid()
                                        table.cell(i + 1, j).fill.fore_color.rgb = RGBColor(0, 255, 0)  # Set background color to green
                                    else:
                                        continue
                                    background_row_counter += 1
                else:
                    for j in range(cols):
                        if j == 1 or j == 3:  # Colored columns
                            for i in range(row_amount):
                                if table_df.iloc[i, j]:
                                    color = color_indexes[background_row_counter]

                                    if color == 'red':
                                        table.cell(i + 1, j).fill.solid()
                                        table.cell(i + 1, j).fill.fore_color.rgb = RGBColor(255, 0, 0)  # Set background color to red
                                    elif color == 'yellow':
                                        table.cell(i + 1, j).fill.solid()
                                        table.cell(i + 1, j).fill.fore_color.rgb = RGBColor(255, 255, 0)  # Set background color to yellow
                                    elif color == 'green':
                                        table.cell(i + 1, j).fill.solid()
                                        table.cell(i + 1, j).fill.fore_color.rgb = RGBColor(0, 255, 0)  # Set background color to green
                                    else:
                                        continue
                                    background_row_counter += 1

            else:
                pass

        # Move slide to the corresponding spot
        last_slide = len(prs.slides)
        move_slide(last_slide - 1, bt.slide_positions[function_index])  # Moving slide into specific location.

        # Updating the slide positions for further appending
        for i in range(function_index, len(bt.slide_positions)):
            bt.slide_positions[i] += 1

        # Updating page amount to end loop
        row_counter += row_amount  # We don't track column values
        page_amount -= 1

    if status_ppt:
        # Creates new pptx name and file if create_ppt function is used first time.
        now = str(datetime.now())  # Taking current time.
        now = now[:-7]
        now = now.replace(':', '.')
        bt.ppt_path = bt.save_path.joinpath('Auswertung_' + now + '.pptx')  # New save path for the appended default slide.
        prs.save(bt.ppt_path.as_posix())  # Changes path to string without backslashes.
    else:
        prs.save(bt.ppt_path.as_posix())  # If second use, we append the used slide instead of default slide
